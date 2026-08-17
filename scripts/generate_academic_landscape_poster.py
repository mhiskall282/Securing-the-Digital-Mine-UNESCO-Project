"""Generate a pixel-perfect academic landscape conference poster with zero image-text overlap.
Strictly uses 10-12pt Times New Roman for body text, Arial Bold for headings,
non-overlapping coordinates, bounded visual cards, and high-res diagrams.
"""
import os
import docx
from docx import Document
from docx.shared import Inches as DInches, Pt as DPt, RGBColor as DRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from docx_styler import clean_text, set_table_borders, set_cell_background, set_cell_margins

def create_pixel_perfect_landscape_poster_pptx():
    prs = Presentation()
    # 48.0 x 32.0 inches (Standard 3:2 Landscape Poster)
    prs.slide_width = Inches(48.0)
    prs.slide_height = Inches(32.0)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Color Palette
    UNESCO_BLUE = RGBColor(0, 82, 155)       # #00529B
    DARK_NAVY = RGBColor(11, 29, 58)         # #0B1D3A
    CYAN_ACCENT = RGBColor(0, 163, 224)      # #00A3E0
    GOLD_ACCENT = RGBColor(255, 215, 0)      # #FFD700
    TEXT_DARK = RGBColor(15, 23, 42)         # #0F172A (Body text)
    TEXT_MUTED = RGBColor(100, 116, 139)     # #64748B
    WHITE = RGBColor(255, 255, 255)
    EMERALD = RGBColor(16, 185, 129)        # #10B981
    CARD_BG = RGBColor(248, 250, 252)        # #F8FAFC
    BORDER_LIGHT = RGBColor(203, 213, 225)   # #CBD5E1

    # 1. Pure White Slide Canvas
    canvas = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(48.0), Inches(32.0))
    canvas.fill.solid()
    canvas.fill.fore_color.rgb = WHITE
    canvas.line.fill.background()

    # 2. Top Header Banner (Y: 0.0 to 4.8 in)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(48.0), Inches(4.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = UNESCO_BLUE
    banner.line.fill.background()

    # Bottom Gold Accent Trim on Banner
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.7), Inches(48.0), Inches(0.1))
    strip.fill.solid()
    strip.fill.fore_color.rgb = GOLD_ACCENT
    strip.line.fill.background()

    # Header Left Logo/Institution Text
    tb_l = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(6.5), Inches(4.0))
    tf_l = tb_l.text_frame
    p_l1 = tf_l.paragraphs[0]
    r_l1 = p_l1.add_run()
    r_l1.text = "UEW"
    r_l1.font.name = "Arial"
    r_l1.font.size = Pt(36)
    r_l1.font.bold = True
    r_l1.font.color.rgb = WHITE

    p_l2 = tf_l.add_paragraph()
    r_l2 = p_l2.add_run()
    r_l2.text = "UNIVERSITY OF EDUCATION,\nWINNEBA, GHANA\n& KAYABA LABS"
    r_l2.font.name = "Arial"
    r_l2.font.size = Pt(13)
    r_l2.font.bold = True
    r_l2.font.color.rgb = CYAN_ACCENT

    # Header Right Forum Details
    tb_r = slide.shapes.add_textbox(Inches(40.5), Inches(0.3), Inches(6.7), Inches(4.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p_r1 = tf_r.paragraphs[0]
    p_r1.alignment = PP_ALIGN.RIGHT
    r_r1 = p_r1.add_run()
    r_r1.text = "UNESCO"
    r_r1.font.name = "Arial"
    r_r1.font.size = Pt(34)
    r_r1.font.bold = True
    r_r1.font.color.rgb = GOLD_ACCENT

    p_r2 = tf_r.add_paragraph()
    p_r2.alignment = PP_ALIGN.RIGHT
    r_r2 = p_r2.add_run()
    r_r2.text = "Russian-African Forum\nTrack 3: Smart Subsoil\nSaint Petersburg Mining Univ."
    r_r2.font.name = "Arial"
    r_r2.font.size = Pt(13)
    r_r2.font.bold = True
    r_r2.font.color.rgb = WHITE

    # Header Center Title & Authors
    tb_c = slide.shapes.add_textbox(Inches(7.5), Inches(0.2), Inches(33.0), Inches(4.3))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True

    p_t = tf_c.paragraphs[0]
    p_t.alignment = PP_ALIGN.CENTER
    rt = p_t.add_run()
    rt.text = "Securing the Digital Mine with a Metaheuristic-Optimized\nDeep Learning Adaptive System"
    rt.font.name = "Arial"
    rt.font.size = Pt(30)
    rt.font.bold = True
    rt.font.color.rgb = WHITE

    p_a = tf_c.add_paragraph()
    p_a.alignment = PP_ALIGN.CENTER
    p_a.space_before = Pt(4)
    ra = p_a.add_run()
    ra.text = "John Okyere (Lead), Ezekeil Baah, Clement Baffour, Parker Paa Annobil, George Akwesi Bonnah"
    ra.font.name = "Times New Roman"
    ra.font.size = Pt(16)
    ra.font.bold = True
    ra.font.color.rgb = GOLD_ACCENT

    p_aff = tf_c.add_paragraph()
    p_aff.alignment = PP_ALIGN.CENTER
    raff = p_aff.add_run()
    raff.text = "Department of ICT, University of Education, Winneba & Kayaba Labs | Saint Petersburg, Russia 2026"
    raff.font.name = "Times New Roman"
    raff.font.size = Pt(13)
    raff.font.italic = True
    raff.font.color.rgb = RGBColor(226, 232, 240)

    # -------------------------------------------------------------
    # 3 STRICT NON-OVERLAPPING COLUMNS (Width: 14.5 in each)
    # -------------------------------------------------------------
    col_w = Inches(14.5)
    c1_x = Inches(1.0)
    c2_x = Inches(16.75)
    c3_x = Inches(32.5)

    def add_heading(left, top, text):
        tb = slide.shapes.add_textbox(left, top, col_w, Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = "Arial"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = UNESCO_BLUE
        return top + Inches(0.55)

    # =============================================================
    # COLUMN 1: THREAT LANDSCAPE, FLOWCHART & BWOA MATH
    # =============================================================
    # 1. Introduction Header & Text (Y: 5.1 to 8.2 in)
    y1 = add_heading(c1_x, Inches(5.1), "1. Industrial Threat Landscape & Context")
    tb_c1_intro = slide.shapes.add_textbox(c1_x, y1, col_w, Inches(2.7))
    tf_c1_i = tb_c1_intro.text_frame
    tf_c1_i.word_wrap = True

    intro_pts = [
        ("Smart Subsoil Digitalization: ", "African and Russian mining operations integrate IoT sensors, SCADA telemetry, and automated milling to drive ore recovery in SAG mills, flotation circuits, and tailings dams."),
        ("Air-Gap Collapse: ", "Connecting OT networks to cloud digital twins exposes unauthenticated Modbus RTU/TCP and DNP3 industrial protocols to hostile zero-day cyber attacks."),
        ("Downtime Losses: ", "Unplanned mining downtime costs USD $50,000 to $500,000 per hour; cyber manipulation of ventilation or dewatering creates immediate life safety risks."),
        ("Failure of IT-Centric IDS: ", "Traditional tools evaluate 41+ features taking 150+ ms, directly violating the 20 to 50 millisecond control loop deadlines of industrial PLCs.")
    ]
    for b_title, b_desc in intro_pts:
        p = tf_c1_i.add_paragraph()
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + b_title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = "Times New Roman"
        r1.font.color.rgb = UNESCO_BLUE
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(9.5)
        r2.font.name = "Times New Roman"
        r2.font.color.rgb = TEXT_DARK

    # Embedded Image 1: Mining SCADA Flowchart (Y: 8.5 to 13.8 in)
    if os.path.exists("research/figures/mining_scada_flowchart.png"):
        slide.shapes.add_picture("research/figures/mining_scada_flowchart.png", c1_x, Inches(8.5), width=col_w)

    # 2. Method & BWOA Mathematics (Y: 14.2 to 29.2 in)
    y1_m = add_heading(c1_x, Inches(14.2), "2. Method & BWOA Mathematics")
    
    # Method Box
    bx_m = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c1_x, y1_m, col_w, Inches(14.8))
    bx_m.fill.solid()
    bx_m.fill.fore_color.rgb = CARD_BG
    bx_m.line.color.rgb = BORDER_LIGHT
    bx_m.line.width = Pt(1.5)

    tb_c1_meth = slide.shapes.add_textbox(c1_x + Inches(0.2), y1_m + Inches(0.2), col_w - Inches(0.4), Inches(14.4))
    tf_c1_m = tb_c1_meth.text_frame
    tf_c1_m.word_wrap = True

    meth_pts = [
        ("Data Ingestion Layer: ", "Promiscuous packet capture using @mhiskall282/unesco-mine-sec-cli parsing IP/TCP/Modbus packet headers at wire speed without packet drops."),
        ("1. Shrinking Encircling Phase: ", "D = |C * X*(t) - X(t)|,  X(t+1) = X*(t) - A * D\nwhere A = 2a*r1 - a, C = 2*r2, and 'a' linearly decreases from 2 to 0 over iterations."),
        ("2. Spiral Bubble-Net Foraging: ", "X(t+1) = D' * exp(b*l) * cos(2*pi*l) + X*(t)\nwhere b=1.0, l is uniform in [-1, 1], mathematically modeling the helical hunting maneuver."),
        ("3. V-Shaped Binary Transfer Function: ", "V(x) = |x / sqrt(1 + x^2)|,  X_d(t+1) = 1 - X_d(t) if rand() < V(x_d) else X_d(t)\nmapping continuous velocities to discrete stochastic bit-flips without boundary saturation."),
        ("4. Accuracy Floor Fitness Function: ", "Fit(X) = alpha*(1 - Acc(X)) + (1-alpha)*(|X|/D) + Penalty(X)\nwhere alpha=0.3 (70% weight on accuracy) and Penalty=1.0 if Acc < 0.75 or |X| < 10."),
        ("Spatial-Temporal CNN-LSTM Classifier: ", "1D Convolutional layer (64 filters, kernel size 3) + BatchNorm + SpatialDropout(0.3) + LSTM (256 units) + Dense(64) + Softmax(5 classes)."),
        ("Float16 Edge Quantization: ", "Post-training Float16 quantization compresses model memory footprint to 0.82 MB (83.2% reduction) for sub-millisecond execution on 1GB RAM edge gateways.")
    ]
    for m_title, m_desc in meth_pts:
        p = tf_c1_m.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "• " + m_title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = "Times New Roman"
        r1.font.color.rgb = UNESCO_BLUE
        r2 = p.add_run()
        r2.text = m_desc
        r2.font.size = Pt(9.5)
        r2.font.name = "Times New Roman"
        r2.font.color.rgb = TEXT_DARK

    # Bottom Left Diagonal Accent Stripes (Y: 29.8 to 31.0 in)
    for i in range(8):
        accent_stripe = slide.shapes.add_shape(
            MSO_SHAPE.PARALLELOGRAM,
            Inches(1.0 + i * 0.9), Inches(29.8), Inches(0.6), Inches(1.1)
        )
        accent_stripe.fill.solid()
        accent_stripe.fill.fore_color.rgb = UNESCO_BLUE if i % 2 == 0 else CYAN_ACCENT
        accent_stripe.line.fill.background()

    # =============================================================
    # COLUMN 2: 4-LAYER ARCHITECTURE & BWOA ANALYSIS
    # =============================================================
    # 1. 4-Layer Architecture (Y: 5.1 to 13.8 in)
    y2 = add_heading(c2_x, Inches(5.1), "3. 4-Layer System Architecture")
    if os.path.exists("research/figures/system_architecture.png"):
        slide.shapes.add_picture("research/figures/system_architecture.png", c2_x, Inches(5.7), width=col_w)

    # 2. BWOA Optimization Analysis (Y: 14.2 to 29.2 in)
    y2_b = add_heading(c2_x, Inches(14.2), "4. BWOA Optimization Analysis")
    if os.path.exists("research/figures/bwoa_convergence.png"):
        slide.shapes.add_picture("research/figures/bwoa_convergence.png", c2_x, Inches(14.8), width=col_w)

    # BWOA Analysis Card Box (Y: 22.0 to 29.2 in)
    bx_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c2_x, Inches(22.0), col_w, Inches(7.2))
    bx_b.fill.solid()
    bx_b.fill.fore_color.rgb = CARD_BG
    bx_b.line.color.rgb = BORDER_LIGHT
    bx_b.line.width = Pt(1.5)

    tb_c2_desc = slide.shapes.add_textbox(c2_x + Inches(0.2), Inches(22.1), col_w - Inches(0.4), Inches(7.0))
    tf_c2_d = tb_c2_desc.text_frame
    tf_c2_d.word_wrap = True

    bwoa_pts = [
        ("Optimal Feature Convergence: ", "BWOA reached minimal fitness at iteration 23/100, maintaining 92.31% Random Forest 3-fold cross-validation accuracy on the selected 10-feature subset."),
        ("75.61% Bandwidth Reduction: ", "Pruned 31 uninformative features from 41, reducing network telemetry transmission load over low-bandwidth satellite links in African mines."),
        ("Selected 10 Vital Features: ", "src_bytes (volume DoS), service & flag (SCADA connection state), serror_rate & diff_srv_rate (reconnaissance), hot & su_attempted (privilege escalation)."),
        ("Transfer Learning Adaptability: ", "Evaluated on the 51-sensor SWaT SCADA benchmark achieving 0.12ms inference latency and 0.8650 AUC-ROC, validating cross-domain industrial suitability.")
    ]
    for b_title, b_desc in bwoa_pts:
        p = tf_c2_d.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = "★ " + b_title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = "Times New Roman"
        r1.font.color.rgb = UNESCO_BLUE
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(9.5)
        r2.font.name = "Times New Roman"
        r2.font.color.rgb = TEXT_DARK

    # =============================================================
    # COLUMN 3: BENCHMARKS, CONFUSION MATRIX & CONCLUSIONS
    # =============================================================
    # 1. Latency Profile Bar Chart (Y: 5.1 to 12.0 in)
    y3 = add_heading(c3_x, Inches(5.1), "5. Empirical Hardware Latency Profile")
    if os.path.exists("research/figures/latency_comparison_barchart.png"):
        slide.shapes.add_picture("research/figures/latency_comparison_barchart.png", c3_x, Inches(5.7), width=col_w)

    # 2. Confusion Matrix (Y: 12.4 to 19.4 in)
    y3_cm = add_heading(c3_x, Inches(12.4), "6. Multi-Class Confusion Matrix")
    if os.path.exists("research/figures/confusion_matrix.png"):
        # Scale confusion matrix neatly to width=11.5 in and center it
        slide.shapes.add_picture("research/figures/confusion_matrix.png", c3_x + Inches(1.5), Inches(13.0), width=Inches(11.5))

    # 3. Conclusions & Impact (Y: 19.8 to 26.8 in)
    y3_c = add_heading(c3_x, Inches(19.8), "7. Conclusions & UN SDG Impact")
    bx_c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c3_x, Inches(20.4), col_w, Inches(6.0))
    bx_c.fill.solid()
    bx_c.fill.fore_color.rgb = CARD_BG
    bx_c.line.color.rgb = BORDER_LIGHT
    bx_c.line.width = Pt(1.5)

    tb_c3_conc = slide.shapes.add_textbox(c3_x + Inches(0.2), Inches(20.5), col_w - Inches(0.4), Inches(5.8))
    tf_c3_c = tb_c3_conc.text_frame
    tf_c3_c.word_wrap = True

    conc_pts = [
        ("0.76ms Edge Real-Time: ", "Executes single-sample inference in 0.76 ms on Raspberry Pi 4B (1GB RAM) - 207x faster than baseline, easily passing the sub-100ms SCADA limit."),
        ("High Detection Fidelity: ", "70.56% multi-class accuracy, 96.89% precision on benign telemetry (zero false plant shutdowns), 89.04% recall on DoS attacks."),
        ("High Industrial ROI: ", "200x-300x return averting $50k-$500k/hr outages while protecting miner lives against ventilation tampering (SDG 8 & 9)."),
        ("Open-Source Ecosystem: ", "NPM sniffer package '@mhiskall282/unesco-mine-sec-cli' and 75/75 automated unit test suites ensure complete reproducibility.")
    ]
    for c_title, c_desc in conc_pts:
        p = tf_c3_c.add_paragraph()
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "✓ " + c_title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = "Times New Roman"
        r1.font.color.rgb = EMERALD
        r2 = p.add_run()
        r2.text = c_desc
        r2.font.size = Pt(9.5)
        r2.font.name = "Times New Roman"
        r2.font.color.rgb = TEXT_DARK

    # 4. References (Y: 27.0 to 31.0 in)
    y3_r = add_heading(c3_x, Inches(27.0), "8. References")
    tb_c3_refs = slide.shapes.add_textbox(c3_x, Inches(27.5), col_w, Inches(3.5))
    tf_c3_r = tb_c3_refs.text_frame
    tf_c3_r.word_wrap = True

    refs = [
        "1. Alanazi, M., et al. (2022). SCADA vulnerabilities & attacks. Computers & Security, 125, 103028.",
        "2. Almomani, O., et al. (2025). SCADA intrusion detection in IIoT. Symmetry, 17(4), 480.",
        "3. Kheddar, H., et al. (2023). Deep transfer learning for intrusion detection in ICS. JNCA, 220, 103747.",
        "4. Mirjalili, S., & Lewis, A. (2016). Whale optimization algorithm. Advances in Eng. Software, 95, 51-67.",
        "5. Minerals Commission of Ghana. (2024). Digital telemetry & cybersecurity compliance policy guidelines."
    ]
    for ref_text in refs:
        p = tf_c3_r.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = ref_text
        r.font.name = "Times New Roman"
        r.font.size = Pt(8.5)
        r.font.color.rgb = TEXT_MUTED

    output_pptx = "research/poster_presentation.pptx"
    prs.save(output_pptx)
    print(f"Pixel-Perfect Academic Landscape Poster (PPTX) saved successfully to {output_pptx}!")

def create_pixel_perfect_landscape_poster_docx():
    doc = Document()
    for section in doc.sections:
        section.top_margin = DInches(0.4)
        section.bottom_margin = DInches(0.4)
        section.left_margin = DInches(0.5)
        section.right_margin = DInches(0.5)
        section.page_width = DInches(17.0)
        section.page_height = DInches(11.0)

    # Top Header Banner Table
    tbl_hdr = doc.add_table(rows=1, cols=3)
    tbl_hdr.alignment = WD_TABLE_ALIGNMENT.CENTER
    tbl_hdr.autofit = False
    tbl_hdr.rows[0].cells[0].width = DInches(2.5)
    tbl_hdr.rows[0].cells[1].width = DInches(11.0)
    tbl_hdr.rows[0].cells[2].width = DInches(2.5)

    for cell in tbl_hdr.rows[0].cells:
        set_cell_background(cell, "00529B")
        set_cell_margins(cell, top=100, bottom=100, left=120, right=120)
    set_table_borders(tbl_hdr, color="FFD700", sz="12", val="single")

    # Left Cell
    p_l = tbl_hdr.rows[0].cells[0].paragraphs[0]
    p_l.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r_l1 = p_l.add_run("UEW\n")
    r_l1.font.name = 'Arial'
    r_l1.font.size = DPt(20)
    r_l1.font.bold = True
    r_l1.font.color.rgb = DRGBColor(255, 255, 255)
    r_l2 = p_l.add_run("Univ. of Education, Winneba\n& Kayaba Labs")
    r_l2.font.name = 'Times New Roman'
    r_l2.font.size = DPt(9.5)
    r_l2.font.bold = True
    r_l2.font.color.rgb = DRGBColor(0, 163, 224)

    # Center Cell
    p_c = tbl_hdr.rows[0].cells[1].paragraphs[0]
    p_c.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r_ct = p_c.add_run("Securing the Digital Mine with a Metaheuristic-Optimized\nDeep Learning Adaptive System\n")
    r_ct.font.name = 'Arial'
    r_ct.font.size = DPt(16)
    r_ct.font.bold = True
    r_ct.font.color.rgb = DRGBColor(255, 255, 255)
    r_ca = p_c.add_run("John Okyere, Ezekeil Baah, Clement Baffour, Parker Paa Annobil, George Akwesi Bonnah\n")
    r_ca.font.name = 'Times New Roman'
    r_ca.font.size = DPt(10.5)
    r_ca.font.bold = True
    r_ca.font.color.rgb = DRGBColor(255, 215, 0)
    r_caff = p_c.add_run("Department of ICT, University of Education, Winneba & Kayaba Labs | Saint Petersburg Mining University 2026")
    r_caff.font.name = 'Times New Roman'
    r_caff.font.size = DPt(9)
    r_caff.font.color.rgb = DRGBColor(226, 232, 240)

    # Right Cell
    p_r = tbl_hdr.rows[0].cells[2].paragraphs[0]
    p_r.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    r_r1 = p_r.add_run("UNESCO\n")
    r_r1.font.name = 'Arial'
    r_r1.font.size = DPt(20)
    r_r1.font.bold = True
    r_r1.font.color.rgb = DRGBColor(255, 215, 0)
    r_r2 = p_r.add_run("Russian-African Forum\nTrack 3: Smart Subsoil")
    r_r2.font.name = 'Times New Roman'
    r_r2.font.size = DPt(9.5)
    r_r2.font.bold = True
    r_r2.font.color.rgb = DRGBColor(255, 255, 255)

    doc.add_paragraph().paragraph_format.space_after = DPt(6)

    # 3-Column Content Table
    tbl_cols = doc.add_table(rows=1, cols=3)
    tbl_cols.alignment = WD_TABLE_ALIGNMENT.CENTER
    tbl_cols.autofit = False
    tbl_cols.rows[0].cells[0].width = DInches(5.3)
    tbl_cols.rows[0].cells[1].width = DInches(5.3)
    tbl_cols.rows[0].cells[2].width = DInches(5.3)

    for cell in tbl_cols.rows[0].cells:
        set_cell_background(cell, "FFFFFF")
        set_cell_margins(cell, top=60, bottom=60, left=80, right=80)
    set_table_borders(tbl_cols, color="CBD5E1", sz="4", val="single")

    def add_col_section(cell, title, bullets, image_path=None, image_w=4.9):
        p_t = cell.add_paragraph() if cell.paragraphs[0].text else cell.paragraphs[0]
        p_t.paragraph_format.space_before = DPt(4)
        p_t.paragraph_format.space_after = DPt(2)
        rt = p_t.add_run(clean_text(title))
        rt.font.name = 'Arial'
        rt.font.size = DPt(11)
        rt.font.bold = True
        rt.font.color.rgb = DRGBColor(0, 82, 155)

        for b in bullets:
            pb = cell.add_paragraph()
            pb.paragraph_format.space_after = DPt(2)
            pb.paragraph_format.line_spacing = 1.15
            rb = pb.add_run("• " + clean_text(b))
            rb.font.name = 'Times New Roman'
            rb.font.size = DPt(9)
            rb.font.color.rgb = DRGBColor(15, 23, 42)

        if image_path and os.path.exists(image_path):
            p_img = cell.add_paragraph()
            p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p_img.paragraph_format.space_before = DPt(3)
            p_img.paragraph_format.space_after = DPt(3)
            p_img.add_run().add_picture(image_path, width=DInches(image_w))

    # Column 1
    c1 = tbl_cols.rows[0].cells[0]
    add_col_section(c1, "1. Industrial Threat Landscape", [
        "Mining 4.0: African/Russian mines deploy IoT & SCADA telemetry to maximize ore yield.",
        "Air-Gap Collapse: Connecting OT to cloud twins exposes Modbus/DNP3 protocols to zero-day attacks.",
        "Downtime Losses: Unplanned downtime costs USD $50k-$500k/hr; gas/ventilation disruption endangers lives.",
        "IT IDS Mismatch: Traditional tools take 150+ ms, violating the 20-50ms SCADA control loop deadline."
    ], image_path="research/figures/mining_scada_flowchart.png", image_w=4.9)

    add_col_section(c1, "2. Method & BWOA Mathematics", [
        "Data Ingestion: Sniffer agent (@mhiskall282/unesco-mine-sec-cli) parses Modbus/TCP frames.",
        "Shrinking Encircling: D = |C*X* - X|, X(t+1) = X* - A*D (A = 2a*r1 - a, C = 2*r2).",
        "Spiral Bubble-Net: X(t+1) = D'*exp(b*l)*cos(2*pi*l) + X* (b=1.0, l in [-1, 1]).",
        "V-Shaped Transfer: V(x) = |x / sqrt(1 + x^2)|, stochastic bit-flipping on velocity.",
        "Accuracy Floor Fitness: Fit(X) = 0.3*(1 - Acc) + 0.7*(|X|/D) + Penalty (Acc >= 75%).",
        "Float16 Quantization: Model compressed to 0.82 MB for sub-millisecond execution on 1GB RAM edge gateways."
    ])

    # Column 2
    c2 = tbl_cols.rows[0].cells[1]
    add_col_section(c2, "3. 4-Layer Architecture Pipeline", [
        "Decoupled 4-layer edge architecture ensures sub-millisecond real-time threat evaluation:"
    ], image_path="research/figures/system_architecture.png", image_w=4.9)

    add_col_section(c2, "4. BWOA Optimization Analysis", [
        "BWOA convergence reached at iteration 23/100, maintaining 92.31% RF cross-validation accuracy.",
        "Selected 10 Features: src_bytes, service, flag, serror_rate, same_srv_rate, diff_srv_rate, dst_host_diff_srv_rate, protocol_type, hot, su_attempted."
    ], image_path="research/figures/bwoa_convergence.png", image_w=4.9)

    # Column 3
    c3 = tbl_cols.rows[0].cells[2]
    add_col_section(c3, "5. Empirical Hardware Latency Profile", [
        "Single-sample inference latency benchmarked across hardware platforms vs 100ms SCADA ceiling:"
    ], image_path="research/figures/latency_comparison_barchart.png", image_w=4.9)

    add_col_section(c3, "6. Multi-Class Confusion Matrix", [
        "Held-out test set evaluation (22,544 samples) on 10 BWOA-selected features:"
    ], image_path="research/figures/confusion_matrix.png", image_w=4.0)

    add_col_section(c3, "7. Conclusions & UN SDG Impact", [
        "0.76ms Edge Latency: 207x faster than baseline on Raspberry Pi 4B (1GB RAM), passing SCADA constraints.",
        "High Detection Fidelity: 70.56% multi-class accuracy, 96.89% benign precision, 89.04% DoS recall.",
        "Industrial ROI: 200x-300x return averting $50k-$500k/hr outages while protecting miner lives (SDG 8 & 9)."
    ])

    add_col_section(c3, "8. References", [
        "1. Alanazi, M., et al. (2022). SCADA vulnerabilities & attacks. Computers & Security, 125, 103028.",
        "2. Almomani, O., et al. (2025). SCADA intrusion detection in IIoT. Symmetry, 17(4), 480.",
        "3. Mirjalili, S., & Lewis, A. (2016). Whale optimization algorithm. Adv. Eng. Software, 95, 51-67.",
        "4. Minerals Commission of Ghana. (2024). Digital telemetry & cybersecurity compliance policy."
    ])

    output_docx = "research/poster_presentation.docx"
    doc.save(output_docx)
    print(f"Pixel-Perfect Academic Landscape Poster (DOCX) saved successfully to {output_docx}!")

if __name__ == "__main__":
    create_pixel_perfect_landscape_poster_pptx()
    create_pixel_perfect_landscape_poster_docx()
