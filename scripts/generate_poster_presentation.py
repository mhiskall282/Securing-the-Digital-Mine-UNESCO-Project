"""Generate high-resolution A0 Portrait Poster Presentation in pptx format."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from docx_styler import clean_text

def create_poster():
    prs = Presentation()
    # A0 Portrait dimensions in inches: 33.11 x 46.81
    prs.slide_width = Inches(33.11)
    prs.slide_height = Inches(46.81)

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Color Palette
    DARK_NAVY = RGBColor(11, 29, 58)
    UNESCO_BLUE = RGBColor(0, 82, 155)
    CYAN = RGBColor(0, 163, 224)
    GOLD = RGBColor(212, 175, 55)
    EMERALD = RGBColor(16, 185, 129)
    WHITE = RGBColor(255, 255, 255)
    TEXT_DARK = RGBColor(30, 41, 59)
    TEXT_MUTED = RGBColor(100, 116, 139)
    BG_CARD = RGBColor(248, 250, 252)
    BORDER_CARD = RGBColor(203, 213, 225)

    # 1. Slide Background (Dark Navy Header + Light Body)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(33.11), Inches(46.81))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
    bg_shape.line.fill.background()

    # 2. Header Banner
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(33.11), Inches(5.8))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_NAVY
    header.line.fill.background()

    # Header Accent Line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.7), Inches(33.11), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    # Title Text
    tb_title = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(31.11), Inches(2.2))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p = tf_title.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "SECURING THE DIGITAL MINE"
    run.font.name = "Arial"
    run.font.size = Pt(56)
    run.font.bold = True
    run.font.color.rgb = GOLD

    p2 = tf_title.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "A Metaheuristic-Optimized Deep Learning Framework for Edge Intrusion Detection in IoT-Enabled Mineral Operations"
    run2.font.name = "Arial"
    run2.font.size = Pt(28)
    run2.font.bold = True
    run2.font.color.rgb = WHITE

    # Subtitle / Forum Details
    tb_sub = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(31.11), Inches(1.2))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    r_sub = p_sub.add_run()
    r_sub.text = "RUSSIAN-AFRICAN FORUM-CONTEST OF YOUNG SCIENTISTS 2026 | TRACK 3: SMART SUBSOIL\nUnder the Auspices of UNESCO | Empress Catherine II Saint Petersburg Mining University"
    r_sub.font.name = "Arial"
    r_sub.font.size = Pt(20)
    r_sub.font.bold = True
    r_sub.font.color.rgb = CYAN

    # Authors
    tb_auth = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(31.11), Inches(1.2))
    tf_auth = tb_auth.text_frame
    tf_auth.word_wrap = True
    p_auth = tf_auth.paragraphs[0]
    p_auth.alignment = PP_ALIGN.CENTER
    r_auth = p_auth.add_run()
    r_auth.text = "John Okyere (Lead) • Ezekeil Baah • Clement Baffour • Parker Paa Annobil • George Akwesi Bonnah\nDepartment of ICT, University of Education, Winneba & Kayaba Labs | hello@johnokyere.xyz"
    r_auth.font.name = "Arial"
    r_auth.font.size = Pt(17)
    r_auth.font.color.rgb = WHITE

    # Helper function for card panels
    def add_card(left, top, width, height, title):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_CARD
        card.line.width = Pt(2)
        
        # Header banner in card
        h_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.1))
        h_card.fill.solid()
        h_card.fill.fore_color.rgb = UNESCO_BLUE
        h_card.line.fill.background()
        
        tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.15), width - Inches(0.6), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.name = "Arial"
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = WHITE
        return card

    # -------------------------------------------------------------
    # PANEL 1: PROBLEM & MOTIVATION (Top Left)
    # -------------------------------------------------------------
    add_card(Inches(1.0), Inches(6.3), Inches(15.0), Inches(11.0), "1. PROBLEM & INDUSTRIAL CONTEXT")
    tb1 = slide.shapes.add_textbox(Inches(1.3), Inches(7.6), Inches(14.4), Inches(9.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    bullets1 = [
        ("Digital Subsoil Transformation: ", "African and Russian mines are rapidly deploying IoT sensors, autonomous haulage, and SCADA telemetry to optimize ore yield."),
        ("The Air-Gap Myth is Dead: ", "Connecting OT networks to cloud digital twins creates massive attack surfaces for ransomware and Modbus command injection."),
        ("Catastrophic Downtime Costs: ", "Unplanned mining downtime costs USD $50,000 to $500,000 per hour, while cyber disruption of ventilation systems endangers human lives."),
        ("Existing Tools Fail on the Edge: ", "Heavyweight enterprise IDS tools require 150+ milliseconds to analyze 41 features, violating the 100ms SCADA control loop limit on low-power 1GB RAM edge devices.")
    ]
    for b_title, b_desc in bullets1:
        p = tf1.add_paragraph()
        p.space_after = Pt(14)
        r1 = p.add_run()
        r1.text = "• " + b_title
        r1.font.bold = True
        r1.font.size = Pt(17)
        r1.font.color.rgb = UNESCO_BLUE
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(16)
        r2.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # PANEL 2: SOLUTION OVERVIEW (Top Right)
    # -------------------------------------------------------------
    add_card(Inches(17.11), Inches(6.3), Inches(15.0), Inches(11.0), "2. 4-LAYER SYSTEM ARCHITECTURE")
    tb2 = slide.shapes.add_textbox(Inches(17.4), Inches(7.6), Inches(14.4), Inches(2.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r = p2.add_run()
    r.text = "A complete edge-to-cloud security pipeline combining metaheuristic feature pruning with spatial-temporal deep learning:"
    r.font.size = Pt(16)
    r.font.color.rgb = TEXT_DARK

    if os.path.exists("research/figures/system_architecture.png"):
        slide.shapes.add_picture("research/figures/system_architecture.png", Inches(17.4), Inches(10.0), width=Inches(14.4))

    # -------------------------------------------------------------
    # PANEL 3: KEY RESULTS & BENCHMARKS (Middle Left)
    # -------------------------------------------------------------
    add_card(Inches(1.0), Inches(17.8), Inches(15.0), Inches(14.0), "3. EXPERIMENTAL BENCHMARK RESULTS")
    tb3 = slide.shapes.add_textbox(Inches(1.3), Inches(19.1), Inches(14.4), Inches(3.2))
    tf3 = tb3.text_frame
    tf3.word_wrap = True

    res_items = [
        ("75.61% Data Pruning: ", "BWOA reduces 41 features down to 10 optimal dimensions while maintaining 92.31% Random Forest CV validation accuracy."),
        ("70.56% Multi-Class Accuracy: ", "Tested on 22,544 held-out NSL-KDD samples with 96.89% precision on Normal traffic and 89.04% recall on DoS attacks."),
        ("0.76ms Edge Latency: ", "Float16 quantized model executes 207x faster than baseline, easily passing the sub-100ms industrial SCADA constraint.")
    ]
    for r_title, r_desc in res_items:
        p = tf3.add_paragraph()
        p.space_after = Pt(10)
        r1 = p.add_run()
        r1.text = "✓ " + r_title
        r1.font.bold = True
        r1.font.size = Pt(17)
        r1.font.color.rgb = EMERALD
        r2 = p.add_run()
        r2.text = r_desc
        r2.font.size = Pt(15.5)
        r2.font.color.rgb = TEXT_DARK

    if os.path.exists("research/figures/latency_comparison_barchart.png"):
        slide.shapes.add_picture("research/figures/latency_comparison_barchart.png", Inches(1.3), Inches(24.2), width=Inches(14.4))

    # -------------------------------------------------------------
    # PANEL 4: DEPLOYMENT READINESS & EDGE SPECS (Middle Right)
    # -------------------------------------------------------------
    add_card(Inches(17.11), Inches(17.8), Inches(15.0), Inches(14.0), "4. EDGE HARDWARE DEPLOYMENT")
    tb4 = slide.shapes.add_textbox(Inches(17.4), Inches(19.1), Inches(14.4), Inches(2.2))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    r = p4.add_run()
    r.text = "Validated on Raspberry Pi 4B (1GB RAM) edge hardware for remote subsoil facilities:"
    r.font.size = Pt(16)
    r.font.color.rgb = TEXT_DARK

    # 4 Stat Boxes
    stats = [
        ("0.76 ms", "Inference Latency\n(Target: < 100 ms)", EMERALD),
        ("0.82 MB", "Model Memory Size\n(83.2% Compression)", UNESCO_BLUE),
        ("290 MB", "Peak RAM Usage\n(Fits 1GB Pi Gateway)", CYAN),
        ("2.5 W", "Power Consumption\n(Solar/Battery Ready)", GOLD)
    ]
    for idx, (val, lbl, col) in enumerate(stats):
        bx = Inches(17.4 + (idx % 2) * 7.3)
        by = Inches(20.4 + (idx // 2) * 2.8)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(7.0), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = col
        box.line.width = Pt(2.5)
        
        t_box = slide.shapes.add_textbox(bx + Inches(0.2), by + Inches(0.2), Inches(6.6), Inches(2.1))
        tf = t_box.text_frame
        p_val = tf.paragraphs[0]
        p_val.alignment = PP_ALIGN.CENTER
        rv = p_val.add_run()
        rv.text = val
        rv.font.name = "Arial"
        rv.font.size = Pt(36)
        rv.font.bold = True
        rv.font.color.rgb = col
        
        p_lbl = tf.add_paragraph()
        p_lbl.alignment = PP_ALIGN.CENTER
        rl = p_lbl.add_run()
        rl.text = lbl
        rl.font.name = "Arial"
        rl.font.size = Pt(14)
        rl.font.color.rgb = TEXT_DARK

    if os.path.exists("research/figures/confusion_matrix.png"):
        slide.shapes.add_picture("research/figures/confusion_matrix.png", Inches(17.4), Inches(26.2), width=Inches(14.4))

    # -------------------------------------------------------------
    # PANEL 5: SUSTAINABILITY & SDG ALIGNMENT (Bottom Left)
    # -------------------------------------------------------------
    add_card(Inches(1.0), Inches(32.3), Inches(15.0), Inches(13.2), "5. ECONOMIC ROI & UN SDG IMPACT")
    tb5 = slide.shapes.add_textbox(Inches(1.3), Inches(33.6), Inches(14.4), Inches(11.5))
    tf5 = tb5.text_frame
    tf5.word_wrap = True

    sdgs = [
        ("SDG 9: Industry, Innovation & Infrastructure", "Hardens critical mining cyber-physical systems against zero-day disruption, safeguarding essential mineral supply chains."),
        ("SDG 8: Decent Work & Economic Growth", "Protects underground worker safety by preventing cyber manipulation of toxic gas sensors and ventilation grids."),
        ("SDG 17: Partnerships for the Goals", "Fosters bilateral Russian-African scientific collaboration, knowledge transfer, and open-source capacity building."),
        ("High Economic ROI (200x - 300x)", "Preventing a single 24-hour ransomware outage on a ball mill or autonomous truck saves $300k to $450k, against an annual deployment cost of under $1,500.")
    ]
    for s_title, s_desc in sdgs:
        p = tf5.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "★ " + s_title + "\n"
        r1.font.bold = True
        r1.font.size = Pt(17)
        r1.font.color.rgb = GOLD
        r2 = p.add_run()
        r2.text = s_desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # PANEL 6: CONCLUSION & REPRODUCIBILITY (Bottom Right)
    # -------------------------------------------------------------
    add_card(Inches(17.11), Inches(32.3), Inches(15.0), Inches(13.2), "6. CONCLUSION & OPEN SOURCE ARTIFACTS")
    tb6 = slide.shapes.add_textbox(Inches(17.4), Inches(33.6), Inches(14.4), Inches(11.5))
    tf6 = tb6.text_frame
    tf6.word_wrap = True

    c_items = [
        ("Complete Open-Source Ecosystem: ", "Full code, training notebooks, benchmarks, and 75 unit test suites are fully open source on GitHub."),
        ("NPM Edge Sniffer Package: ", "Install instantly on any gateway via '@mhiskall282/unesco-mine-sec-cli' from GitHub Packages."),
        ("UNESCO Forum Delegation: ", "Presented by the University of Education, Winneba & Kayaba Labs research team at Empress Catherine II Saint Petersburg Mining University."),
        ("Repository URL: ", "https://github.com/mhiskall282/Securing-the-Digital-Mine-UNESCO-Project"),
        ("Correspondence Email: ", "hello@johnokyere.xyz | Portfolio: https://johnokyere.xyz")
    ]
    for c_title, c_desc in c_items:
        p = tf6.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "► " + c_title
        r1.font.bold = True
        r1.font.size = Pt(17)
        r1.font.color.rgb = UNESCO_BLUE
        r2 = p.add_run()
        r2.text = c_desc
        r2.font.size = Pt(15.5)
        r2.font.color.rgb = TEXT_DARK

    # Save presentation
    output_path = "research/poster_presentation.pptx"
    prs.save(output_path)
    print(f"Poster Presentation saved successfully to {output_path}!")

if __name__ == "__main__":
    create_poster()
